#!/usr/bin/env python3
"""Add native PowerPoint slide-number fields to an image-deck PPTX."""
from __future__ import annotations

import argparse
from io import BytesIO
import re
import uuid
from pathlib import Path

from PIL import Image, ImageStat
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.util import Inches

SHAPE_NAME = "image-deck-native-slide-number"
EMU_PER_INCH = 914400


def parse_color(value: str) -> RGBColor | None:
    if value.strip().lower() == "auto":
        return None
    value = value.strip().lstrip("#")
    if not re.fullmatch(r"[0-9A-Fa-f]{6}", value):
        raise argparse.ArgumentTypeError("color must be 'auto' or a 6-digit RRGGBB value")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def choose_contrast_color(slide, prs, *, right_in: float, bottom_in: float,
                          width_in: float, height_in: float) -> RGBColor:
    pictures = [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
    if not pictures:
        return RGBColor(255, 255, 255)

    picture = max(pictures, key=lambda shape: shape.width * shape.height)
    try:
        with Image.open(BytesIO(picture.image.blob)).convert("RGB") as image:
            slide_w = prs.slide_width / EMU_PER_INCH
            slide_h = prs.slide_height / EMU_PER_INCH
            left_ratio = max(0.0, (slide_w - right_in - width_in) / slide_w)
            top_ratio = max(0.0, (slide_h - bottom_in - height_in) / slide_h)
            right_ratio = min(1.0, (slide_w - right_in) / slide_w)
            bottom_ratio = min(1.0, (slide_h - bottom_in) / slide_h)
            crop = image.crop((
                int(image.width * left_ratio),
                int(image.height * top_ratio),
                max(1, int(image.width * right_ratio)),
                max(1, int(image.height * bottom_ratio)),
            ))
            red, green, blue = ImageStat.Stat(crop).mean[:3]
    except Exception:
        return RGBColor(255, 255, 255)

    luminance = (0.2126 * red + 0.7152 * green + 0.0722 * blue) / 255
    return RGBColor(0, 0, 0) if luminance >= 0.58 else RGBColor(255, 255, 255)


def add_slide_number(slide, slide_number: int, prs, *, font_name: str,
                     font_size_pt: float, color: RGBColor, right_in: float,
                     bottom_in: float, width_in: float, height_in: float,
                     placeholder_idx: str) -> None:
    slide_w = prs.slide_width / EMU_PER_INCH
    slide_h = prs.slide_height / EMU_PER_INCH
    left = max(0.0, slide_w - right_in - width_in)
    top = max(0.0, slide_h - bottom_in - height_in)

    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width_in), Inches(height_in)
    )
    box.name = SHAPE_NAME
    box.fill.background()
    box.line.fill.background()

    # PowerPoint only renders a slidenum field reliably when the field shape is
    # also a real slide-number placeholder.  LibreOffice and WPS accept a field
    # inside an ordinary text box, which previously hid this compatibility gap.
    c_nv_sp_pr = box._element.nvSpPr.cNvSpPr
    c_nv_sp_pr.attrib.pop("txBox", None)
    c_nv_sp_pr.append(parse_xml(
        f'<a:spLocks {nsdecls("a")} noGrp="1"/>'
    ))
    nv_pr = box._element.nvSpPr.nvPr
    nv_pr.append(parse_xml(
        f'<p:ph {nsdecls("p")} type="sldNum" sz="quarter" '
        f'idx="{placeholder_idx}"/>'
    ))

    text_frame = box.text_frame
    text_frame.clear()
    text_frame.word_wrap = False
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0

    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.RIGHT
    rgb = "".join(f"{component:02X}" for component in color)
    field_xml = (
        f'<a:fld {nsdecls("a")} id="{{{str(uuid.uuid4()).upper()}}}" type="slidenum">'
        f'<a:rPr lang="en-US" sz="{int(round(font_size_pt * 100))}">'
        f'<a:latin typeface="{font_name}"/>'
        f'<a:solidFill><a:srgbClr val="{rgb}"/></a:solidFill>'
        f'</a:rPr><a:t>{slide_number}</a:t></a:fld>'
    )
    paragraph._p.append(parse_xml(field_xml))
    paragraph._p.append(parse_xml(
        f'<a:endParaRPr {nsdecls("a")} lang="en-US"/>'
    ))


def slide_number_placeholder_idx(slide) -> str:
    layout_fields = slide.slide_layout._element.xpath(
        './/p:ph[@type="sldNum"]'
    )
    master_fields = slide.slide_layout.slide_master._element.xpath(
        './/p:ph[@type="sldNum"]'
    )
    if len(layout_fields) != 1 or len(master_fields) != 1:
        raise RuntimeError(
            "PowerPoint-compatible native numbering requires exactly one "
            "sldNum placeholder in both the slide layout and slide master"
        )
    return layout_fields[0].get("idx", "0")


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("input", type=Path)
    parser.add_argument("output", type=Path)
    parser.add_argument("--font", default="Arial")
    parser.add_argument("--size", type=float, default=16.0)
    parser.add_argument("--color", type=parse_color, default=None,
                        help="'auto' (default) or a 6-digit RRGGBB value")
    parser.add_argument("--right", type=float, default=0.28)
    parser.add_argument("--bottom", type=float, default=0.18)
    parser.add_argument("--width", type=float, default=0.72)
    parser.add_argument("--height", type=float, default=0.44)
    parser.add_argument("--include-cover", action="store_true")
    args = parser.parse_args()

    presentation = Presentation(args.input)

    for slide in presentation.slides:
        for shape in list(slide.shapes):
            if shape.name == SHAPE_NAME:
                shape._element.getparent().remove(shape._element)

    for slide_number, slide in enumerate(presentation.slides, start=1):
        if slide_number == 1 and not args.include_cover:
            continue
        color = args.color or choose_contrast_color(
            slide,
            presentation,
            right_in=args.right,
            bottom_in=args.bottom,
            width_in=args.width,
            height_in=args.height,
        )
        add_slide_number(
            slide,
            slide_number,
            presentation,
            font_name=args.font,
            font_size_pt=args.size,
            color=color,
            right_in=args.right,
            bottom_in=args.bottom,
            width_in=args.width,
            height_in=args.height,
            placeholder_idx=slide_number_placeholder_idx(slide),
        )

    args.output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(args.output)


if __name__ == "__main__":
    main()
