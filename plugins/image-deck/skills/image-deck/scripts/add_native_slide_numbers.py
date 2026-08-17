\
#!/usr/bin/env python3
"""Add consistent native PowerPoint slide-number fields to an image-deck PPTX."""
from __future__ import annotations

import argparse
import re
import uuid
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.util import Inches

SHAPE_NAME = "image-deck-native-slide-number"
EMU_PER_INCH = 914400


def parse_color(value: str) -> RGBColor:
    value = value.strip().lstrip("#")
    if not re.fullmatch(r"[0-9A-Fa-f]{6}", value):
        raise argparse.ArgumentTypeError("color must be a 6-digit RRGGBB value")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def add_slide_number(slide, slide_number: int, prs, *, font_name: str, font_size_pt: float,
                     color: RGBColor, right_in: float, bottom_in: float,
                     width_in: float, height_in: float) -> None:
    slide_w = prs.slide_width / EMU_PER_INCH
    slide_h = prs.slide_height / EMU_PER_INCH
    left = max(0.0, slide_w - right_in - width_in)
    top = max(0.0, slide_h - bottom_in - height_in)

    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width_in), Inches(height_in)
    )
    box.name = SHAPE_NAME
    box.fill.background()
    box.line.fill.background()

    tf = box.text_frame
    tf.clear()
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    rgb = "".join(f"{component:02X}" for component in color)
    field_xml = (
        f'<a:fld {nsdecls("a")} id="{{{uuid.uuid4().hex.upper()}}}" type="slidenum">'
        f'<a:rPr lang="en-US" sz="{int(round(font_size_pt * 100))}">'
        f'<a:latin typeface="{font_name}"/>'
        f'<a:solidFill><a:srgbClr val="{rgb}"/></a:solidFill>'
        f'</a:rPr><a:t>{slide_number}</a:t></a:fld>'
    )
    p._p.append(parse_xml(field_xml))


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("input", type=Path)
    parser.add_argument("output", type=Path)
    parser.add_argument("--font", default="Arial")
    parser.add_argument("--size", type=float, default=9.0)
    parser.add_argument("--color", type=parse_color, default=parse_color("FFFFFF"))
    parser.add_argument("--right", type=float, default=0.28)
    parser.add_argument("--bottom", type=float, default=0.18)
    parser.add_argument("--width", type=float, default=0.55)
    parser.add_argument("--height", type=float, default=0.24)
    parser.add_argument("--include-cover", action="store_true")
    args = parser.parse_args()

    prs = Presentation(args.input)

    for slide in prs.slides:
        for shape in list(slide.shapes):
            if shape.name == SHAPE_NAME:
                shape._element.getparent().remove(shape._element)

    for slide_number, slide in enumerate(prs.slides, start=1):
        if slide_number == 1 and not args.include_cover:
            continue
        add_slide_number(
            slide, slide_number, prs,
            font_name=args.font,
            font_size_pt=args.size,
            color=args.color,
            right_in=args.right,
            bottom_in=args.bottom,
            width_in=args.width,
            height_in=args.height,
        )

    args.output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(args.output)


if __name__ == "__main__":
    main()
