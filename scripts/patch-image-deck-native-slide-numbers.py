#!/usr/bin/env python3
from __future__ import annotations

import json
from pathlib import Path
from textwrap import dedent

ROOT = Path(__file__).resolve().parents[1]
SKILL_PATHS = [
    ROOT / "skills/image-deck/SKILL.md",
    ROOT / "plugins/image-deck/skills/image-deck/SKILL.md",
]
REFERENCE_PATHS = [
    ROOT / "skills/image-deck/references/prompt-patterns.md",
    ROOT / "plugins/image-deck/skills/image-deck/references/prompt-patterns.md",
]

HELPER = dedent(r'''\
#!/usr/bin/env python3
"""Add consistent native PowerPoint slide-number fields to an image-deck PPTX."""
from __future__ import annotations

import argparse
import re
import uuid
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.util import Inches

SHAPE_NAME = "image-deck-native-slide-number"
EMU_PER_INCH = 914400


def parse_color(value: str) -> RGBColor:
    value = value.strip().lstrip("#")
    if not re.fullmatch(r"[0-9A-Fa-f]{6}", value):
        raise argparse.ArgumentTypeError("color must be a 6-digit RRGGBB value")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def add_slide_number(slide, slide_number: int, prs, *, font_name: str, font_size_pt: float,
                     color: RGBColor, right_in: float, bottom_in: float,
                     width_in: float, height_in: float) -> None:
    slide_w = prs.slide_width / EMU_PER_INCH
    slide_h = prs.slide_height / EMU_PER_INCH
    left = max(0.0, slide_w - right_in - width_in)
    top = max(0.0, slide_h - bottom_in - height_in)

    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width_in), Inches(height_in)
    )
    box.name = SHAPE_NAME
    box.fill.background()
    box.line.fill.background()

    tf = box.text_frame
    tf.clear()
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    rgb = "".join(f"{component:02X}" for component in color)
    field_xml = (
        f'<a:fld {nsdecls("a")} id="{{{uuid.uuid4().hex.upper()}}}" type="slidenum">'
        f'<a:rPr lang="en-US" sz="{int(round(font_size_pt * 100))}">'
        f'<a:latin typeface="{font_name}"/>'
        f'<a:solidFill><a:srgbClr val="{rgb}"/></a:solidFill>'
        f'</a:rPr><a:t>{slide_number}</a:t></a:fld>'
    )
    p._p.append(parse_xml(field_xml))


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("input", type=Path)
    parser.add_argument("output", type=Path)
    parser.add_argument("--font", default="Arial")
    parser.add_argument("--size", type=float, default=9.0)
    parser.add_argument("--color", type=parse_color, default=parse_color("FFFFFF"))
    parser.add_argument("--right", type=float, default=0.28)
    parser.add_argument("--bottom", type=float, default=0.18)
    parser.add_argument("--width", type=float, default=0.55)
    parser.add_argument("--height", type=float, default=0.24)
    parser.add_argument("--include-cover", action="store_true")
    args = parser.parse_args()

    prs = Presentation(args.input)

    for slide in prs.slides:
        for shape in list(slide.shapes):
            if shape.name == SHAPE_NAME:
                shape._element.getparent().remove(shape._element)

    for slide_number, slide in enumerate(prs.slides, start=1):
        if slide_number == 1 and not args.include_cover:
            continue
        add_slide_number(
            slide, slide_number, prs,
            font_name=args.font,
            font_size_pt=args.size,
            color=args.color,
            right_in=args.right,
            bottom_in=args.bottom,
            width_in=args.width,
            height_in=args.height,
        )

    args.output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(args.output)


if __name__ == "__main__":
    main()
''')

NATIVE_SECTION = dedent('''\
### Native PowerPoint slide numbering

Page numbers are document-level structure, not generated artwork.

- Never generate a page number, slide number, or footer number inside the slide image.
- Cover slide: no page number by default.
- Inner slides: add exactly one native PowerPoint slide-number field during PPTX assembly.
- The field displays the actual PowerPoint slide number, so the second slide displays `2` when the cover is slide 1.
- Default format: plain integer (`2`, `3`, `4`…), with no leading zero and no `x / total` suffix.
- Default position: bottom-right, with a fixed 0.28-inch right inset and 0.18-inch bottom inset on a 16:9 slide.
- Default font: Arial, 9 pt, right-aligned. Font and geometry stay fixed across the deck.
- Default color: choose black or white from the approved visual system based on footer contrast; do not change position or font from slide to slide.
- Reserve a quiet, low-detail footer-safe area in the generated image so the native field remains legible.
- Use `scripts/add_native_slide_numbers.py` for the final PPTX assembly step. It injects a real PowerPoint `slidenum` field, not a static text value.
- If slides are inserted, removed, or reordered after assembly, rerun the helper so the field exists on every non-cover slide. The field itself remains dynamic and updates to the actual slide number in PowerPoint.
- The native page number is the only permitted post-generation visible text element. Do not add other text overlays to compensate for generated-image text problems.
''')


def replace_once(text: str, old: str, new: str, label: str) -> str:
    count = text.count(old)
    if count != 1:
        raise RuntimeError(f"{label}: expected 1 occurrence, found {count}")
    return text.replace(old, new, 1)


def patch_skill(path: Path) -> None:
    text = path.read_text(encoding="utf-8")

    text = replace_once(
        text,
        "If the slide needs a title, caption, number, label, chart title, or short bullet, that visible text must be requested in the Codex `image_gen` (GPT Image 2) prompt and must appear inside the generated image itself.\n\nWhen assembling PPTX/PDF, each slide must contain the generated image as the only visible slide content. Do not add separate text boxes, captions, page numbers, icons, shapes, charts, or labels after generation. If text is missing, wrong, or unreadable, regenerate the slide image instead of overlaying corrected text locally.",
        "If the slide needs a title, caption, label, chart title, or short bullet, that visible text must be requested in the Codex `image_gen` (GPT Image 2) prompt and must appear inside the generated image itself. **Slide numbers are the one deliberate exception:** never generate, draw, or request a page number inside the slide image. Page numbers are added after image generation as native PowerPoint slide-number fields.\n\nWhen assembling PPTX/PDF, each slide must contain the generated image as the only generated slide content. For inner slides, the assembly step may add exactly one native PowerPoint slide-number field in the standardized footer position. Do not add any other separate text boxes, captions, labels, icons, shapes, or charts after generation. If generated slide text is missing, wrong, or unreadable, regenerate the slide image instead of overlaying corrected text locally. If the native slide number is wrong or misplaced, fix the PowerPoint field/style at assembly time rather than regenerating the image.",
        f"{path}: generation boundary",
    )

    text = replace_once(
        text,
        "- Every slide's visible title, explanatory text, labels, bullets/callouts, and short copy, when used, are inside the generated image itself, not overlaid later.",
        "- Every slide's generated title, explanatory text, labels, bullets/callouts, and short copy, when used, are inside the generated image itself, not overlaid later. Native slide numbers are excluded from this rule and are added programmatically as PowerPoint slide-number fields.",
        f"{path}: done criteria text",
    )

    text = replace_once(
        text,
        "- If the final deck is assembled, verify that each slide is one full-bleed generated image with no extra visible text/shape objects.",
        "- If the final deck is assembled, verify that each slide is one full-bleed generated image plus, on inner slides, one standardized native PowerPoint slide-number field; no other visible text/shape objects may be added.",
        f"{path}: done criteria assembly",
    )

    text = replace_once(
        text,
        "All visible content must be generated inside each slide image. This includes the title, subtitle, bullet-style callouts, labels, page number, section tag, captions, and any short copy.",
        "All generated slide content must be generated inside each slide image. This includes the title, subtitle, bullet-style callouts, labels, section tag, captions, and any short copy. **Do not include the page number in the generated image.** The page number is a native PowerPoint field added during final PPTX assembly.",
        f"{path}: full PPT page mode",
    )

    text = replace_once(
        text,
        "- grid, title zone, text/callout zone, main visual zone, footer/page marker zone for inner pages",
        "- grid, title zone, text/callout zone, main visual zone, footer-safe zone reserved for the native PowerPoint slide-number field on inner pages",
        f"{path}: visual bible footer",
    )

    text = replace_once(
        text,
        "- **Add new slides:** extend the deck spine, create prompts using the same visual bible and the same prompt-group format, ask for prompt approval for the new slides, generate them through Codex `image_gen` (GPT Image 2), insert them into the deck, and update numbering if the deck uses generated page markers.",
        "- **Add new slides:** extend the deck spine, create prompts using the same visual bible and the same prompt-group format, ask for prompt approval for the new slides, generate them through Codex `image_gen` (GPT Image 2), insert them into the deck, and rerun the native slide-number helper so every non-cover slide has the same field/style.",
        f"{path}: add slides numbering",
    )

    text = replace_once(
        text,
        "Do not repair a delivered deck by adding PPT text boxes or shapes. Any visible change must be made by regenerating the affected full-slide image.",
        "Do not repair a delivered deck by adding PPT text boxes or shapes. Any visible content change must be made by regenerating the affected full-slide image. The only permitted post-generation visible element is the standardized native PowerPoint slide-number field.",
        f"{path}: post-delivery repair rule",
    )

    text = replace_once(
        text,
        "- the original aspect ratio, palette, lighting, margins, typography mood, page markers, and role system",
        "- the original aspect ratio, palette, lighting, margins, typography mood, footer-safe area, native slide-number field style, and role system",
        f"{path}: revision consistency",
    )

    text = replace_once(
        text,
        "- Prefer recurring page devices: corner number, chapter tag, consistent title position, repeated frame/grid.",
        "- Prefer recurring page devices: chapter tag, consistent title position, repeated frame/grid, and a reserved footer-safe zone for the native slide number.",
        f"{path}: practical page devices",
    )

    text = text.replace("page markers", "native slide-number fields")
    text = text.replace("page marker", "native slide-number field")

    anchor = "### 12. Revise or expand after delivery"
    if NATIVE_SECTION.strip() not in text:
        if text.count(anchor) != 1:
            raise RuntimeError(f"{path}: native numbering insertion anchor not found exactly once")
        text = text.replace(anchor, NATIVE_SECTION + "\n" + anchor, 1)

    handoff_anchor = "- any residual limitation, especially generated text accuracy inside images"
    text = replace_once(
        text,
        handoff_anchor,
        handoff_anchor + "\n- native slide-number field verification: cover excluded, inner slides numbered, standardized font/position/style",
        f"{path}: handoff checklist",
    )

    path.write_text(text, encoding="utf-8")


def patch_references(path: Path) -> None:
    text = path.read_text(encoding="utf-8")

    replacements = [
        (
            "Inner-page layout grid: <role-appropriate title/claim zone>, <text/callout zone>, <main visual zone>, <footer/page marker zone>, generous safe margins.",
            "Inner-page layout grid: <role-appropriate title/claim zone>, <text/callout zone>, <main visual zone>, <footer-safe zone reserved for the native PowerPoint slide-number field>, generous safe margins.",
        ),
        (
            "Text rule: all visible text must be generated inside this image; do not leave blank title areas for later editing. Match text density to slide role and selected content-density mode.",
            "Text rule: all generated slide content text must be generated inside this image; do not leave blank title areas for later editing. Do not generate the page number in the image; it is added later as a native PowerPoint slide-number field. Match text density to slide role and selected content-density mode.",
        ),
        (
            "Composition: <where generated text belongs, where the main visual goes, how labels/callouts attach to the visual, repeated generated footer/page marker if needed. For cover slides, use only the main title plus optional subtitle unless the user explicitly asked for more cover text. For content slides, allow naturally generated supporting detail text when it makes the slide richer and more complete. Keep style consistent without copying another slide's background>",
            "Composition: <where generated text belongs, where the main visual goes, how labels/callouts attach to the visual, and where to leave a clean footer-safe zone for the native PowerPoint slide-number field. Never render the page number inside the image. For cover slides, use only the main title plus optional subtitle unless the user explicitly asked for more cover text. For content slides, allow naturally generated supporting detail text when it makes the slide richer and more complete. Keep style consistent without copying another slide's background>",
        ),
        (
            "- <list exact slide text here: title/claim, captions, bullets, callouts, labels, section tag, or page marker as appropriate to this slide role>",
            "- <list exact generated slide text here: title/claim, captions, bullets, callouts, labels, or section tag as appropriate to this slide role; do not include the page number>",
        ),
        (
            "Consistency rule: keep the same palette, lighting, margin, typography mood, page marker, and graphic language as the approved visual bible and other prompt groups.",
            "Consistency rule: keep the same palette, lighting, margin, typography mood, footer-safe area, native slide-number field placement, and graphic language as the approved visual bible and other prompt groups.",
        ),
        (
            "Regenerate slide <N> with the same locked visual bible and same content. Fix only this issue: <issue>.\nKeep the title area cleaner, reduce decorative clutter, and preserve the same palette and page marker style as the approved master sample.",
            "Regenerate slide <N> with the same locked visual bible and same content. Fix only this issue: <issue>.\nKeep the title area cleaner, reduce decorative clutter, and preserve the same palette and footer-safe area as the approved master sample. The native page number is added by the PPTX assembly step, not by image generation.",
        ),
    ]

    for old, new in replacements:
        if old not in text:
            raise RuntimeError(f"{path}: reference phrase not found: {old[:80]!r}")
        text = text.replace(old, new, 1)

    text = text.replace("page marker", "native slide-number field")
    text = text.replace("page markers", "native slide-number fields")

    if "Do not generate the page number in the image" not in text:
        raise RuntimeError(f"{path}: native page-number prohibition missing")

    path.write_text(text, encoding="utf-8")


def main() -> None:
    for path in SKILL_PATHS:
        patch_skill(path)
    for path in REFERENCE_PATHS:
        patch_references(path)

    for skill_dir in [ROOT / "skills/image-deck", ROOT / "plugins/image-deck/skills/image-deck"]:
        helper = skill_dir / "scripts/add_native_slide_numbers.py"
        helper.parent.mkdir(parents=True, exist_ok=True)
        helper.write_text(HELPER, encoding="utf-8")
        helper.chmod(0o755)

    plugin_json_path = ROOT / "plugins/image-deck/.codex-plugin/plugin.json"
    plugin = json.loads(plugin_json_path.read_text(encoding="utf-8"))
    if plugin.get("version") == "1.0.0":
        plugin["version"] = "1.0.1"
    plugin["interface"]["longDescription"] = (
        "Create polished visual decks with Codex built-in image generation. "
        "image-deck researches or extracts source material, designs the narrative and visual system, "
        "shows prompts for review, generates one master sample for approval, and then produces the remaining "
        "full-slide images for optional PPTX or PDF assembly. PPTX output uses standardized native PowerPoint "
        "slide-number fields rather than generated page-number artwork."
    )
    plugin_json_path.write_text(json.dumps(plugin, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")

    publish_workflow = ROOT / ".github/workflows/publish-image-deck.yml"
    publish_workflow.write_text(dedent('''\
name: Publish image-deck to ClawHub

on:
  push:
    branches: [main]
    paths:
      - 'skills/image-deck/**'
      - '.github/workflows/publish-image-deck.yml'
  workflow_dispatch:

permissions:
  contents: read

jobs:
  publish:
    uses: openclaw/clawhub/.github/workflows/skill-publish.yml@v1
    with:
      owner: tseng71
      skill_path: skills/image-deck
      dry_run: false
    secrets:
      clawhub_token: ${{ secrets.CLAWHUB_TOKEN }}
'''), encoding="utf-8")

    local_sync = ROOT / "scripts/sync-image-deck-local.sh"
    local_sync.write_text(dedent('''\
#!/usr/bin/env bash
set -euo pipefail

repo_root="$(cd "$(dirname "${BASH_SOURCE[0]}")/.." && pwd)"
source_skill="$repo_root/skills/image-deck"
plugin_skill="$repo_root/plugins/image-deck/skills/image-deck"
target_dir="${CODEX_IMAGE_DECK_DIR:-$HOME/.codex/skills/image-deck}"

if [[ ! -d "$source_skill" ]]; then
  echo "Missing source skill: $source_skill" >&2
  exit 1
fi

rm -rf "$target_dir"
mkdir -p "$target_dir"
cp -R "$source_skill/." "$target_dir/"

diff -qr "$source_skill" "$plugin_skill"
echo "Synced image-deck to $target_dir"
'''), encoding="utf-8")
    local_sync.chmod(0o755)

    print("image-deck native slide-number migration prepared")


if __name__ == "__main__":
    main()
